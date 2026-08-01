from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Anchor all paths to the project ROOT (the folder above /notebooks).
# Works no matter which folder you run the script from.
ROOT  = Path(__file__).resolve().parent.parent
DELIV = ROOT / "deliverables"
DELIV.mkdir(exist_ok=True)

# ---------- palette & type ----------
INK  = RGBColor(0x1A, 0x1A, 0x2E)
GREY = RGBColor(0x55, 0x60, 0x6E)
FONT = "Segoe UI"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def _fmt(run, size, color, bold=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold

def title_slide(title, subtitle_lines):
    s = prs.slides.add_slide(BLANK)
    tb = s.shapes.add_textbox(Inches(0.9), Inches(2.6), Inches(11.5), Inches(1.6))
    tb.text_frame.word_wrap = True
    r = tb.text_frame.paragraphs[0].add_run(); r.text = title
    _fmt(r, 40, INK, bold=True)
    sb = s.shapes.add_textbox(Inches(0.9), Inches(4.4), Inches(11.5), Inches(1.5))
    sb.text_frame.word_wrap = True
    for i, line in enumerate(subtitle_lines):
        p = sb.text_frame.paragraphs[0] if i == 0 else sb.text_frame.add_paragraph()
        r = p.add_run(); r.text = line; _fmt(r, 16, GREY); p.space_after = Pt(4)

def content_slide(action_title, bullets):
    s = prs.slides.add_slide(BLANK)
    tb = s.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(1.4))
    tb.text_frame.word_wrap = True
    r = tb.text_frame.paragraphs[0].add_run(); r.text = action_title
    _fmt(r, 30, INK, bold=True)
    bb = s.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(12), Inches(4.8))
    bb.text_frame.word_wrap = True
    for i, b in enumerate(bullets):
        p = bb.text_frame.paragraphs[0] if i == 0 else bb.text_frame.add_paragraph()
        r = p.add_run(); r.text = "\u2022  " + b; _fmt(r, 18, INK)
        p.space_after = Pt(12); p.line_spacing = 1.15

def dashboard_slide(action_title, image_path, caption=None):
    s = prs.slides.add_slide(BLANK)
    tb = s.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(1.0))
    tb.text_frame.word_wrap = True
    r = tb.text_frame.paragraphs[0].add_run(); r.text = action_title
    _fmt(r, 28, INK, bold=True)
    s.shapes.add_picture(str(image_path), Inches(0.9), Inches(1.7), width=Inches(11.5))
    if caption:
        cb = s.shapes.add_textbox(Inches(0.7), Inches(6.95), Inches(12), Inches(0.4))
        r = cb.text_frame.paragraphs[0].add_run(); r.text = caption
        _fmt(r, 12, GREY)

# ==================== YOUR DECK ====================
title_slide(
    "CRM Sales Pipeline: Where Winnable Revenue Leaks",
    ["Prepared by Rajkumar Vijayan",
     "Client: MavenTech  \u00b7  Analysis of 8,800 sales opportunities",
     "August 2026"])

content_slide(
    "MavenTech can recover at least $237K in won revenue from sales coaching alone",
    ["Situation: MavenTech's sales team wins 63.2% of closed deals and $10.0M in revenue, but has no visibility into where value is lost.",
     "Complication: A diagnosis of 8,800 opportunities reveals two recoverable gaps.",
     "Findings: a 15-point win-rate spread across agents, and $4.9M in losses concentrated in three products.",
     "Recommendation: coaching below-median agents to the team average alone represents at least $237K in recoverable revenue, within one quarter."])

content_slide(
    "MavenTech's new CRM tracks deals but leaves leadership blind to where they're won or lost",
    ["A new CRM system captures every opportunity, but the data has no visibility outside the platform.",
     "Leadership cannot see which agents, products, or stages drive won vs. lost revenue.",
     "Question: where is the team losing winnable revenue, and what is it worth to fix?"])

content_slide(
    "We traced 8,800 opportunities end-to-end, from raw CRM extract to validated findings",
    ["Data: 8,800 B2B deals across four linked tables (opportunities, agents, accounts, products).",
     "Method: cleaned and validated in Python, diagnosed the funnel, cross-checked every figure in Tableau.",
     "Examined three levers: agent consistency, lost-deal concentration, and discount leakage."])

dashboard_slide(
    "One view of the whole pipeline: win rate, revenue, and where deals are lost",
    DELIV / "CRM_Sales_Dashboard.png",
    "Live, interactive dashboard published on Tableau Public.")

content_slide(
    "The team wins 63% of deals \u2014 but individual win rates swing 15 points",
    ["Median agent win rate is 63.6%, but rates range from 55% to 70% across the team.",
     "The spread holds on the same products and accounts \u2014 so it reflects execution, not territory.",
     "That gap is coachable, which makes it the most actionable lever."])

content_slide(
    "Lifting below-median agents to the team average is worth at least $237K",
    ["Counting only agents below median, lifted only to the median (a deliberately conservative floor).",
     "That yields 101 additional wins at an average deal value of $2,361.",
     "Estimated upside: $237,433 \u2014 a 2.4% lift on current won revenue.",
     "This is the floor: lifting to top-quartile performance would recover materially more."])

content_slide(
    "Three products drive ~$4.9M of lost revenue \u2014 and GTX Pro loses the most per deal",
    ["GTX Pro ($2.0M), MG Advanced ($1.46M), and GTX Plus Pro ($1.46M) account for most losses.",
     "GTX Pro loses fewer deals than MG Advanced but costs more per loss \u2014 it is higher-value.",
     "Prioritising GTX Pro win-back therefore returns more per deal recovered."])

content_slide(
    "A single CRM naming inconsistency was hiding the #1 lost-revenue product",
    ["The CRM stored the same product as both 'GTXPro' and 'GTX Pro', silently breaking table joins.",
     "This undercounted GTX Pro's losses until the naming was standardised.",
     "Fixing one data-quality issue changed which product looked worst \u2014 a process gap worth closing."])

content_slide(
    "Three targeted actions capture the opportunity",
    ["Coach bottom-quartile agents toward the median team's playbook and cadence.",
     "Prioritise GTX Pro win-back, given its high revenue-per-lost-deal.",
     "Enforce consistent product naming in the CRM to keep reporting reliable."])

content_slide(
    "At least $237K is recoverable from coaching alone, measurable within one quarter",
    ["Track win rate by agent monthly and close the gap to the team median.",
     "Run a targeted GTX Pro loss review to recover high-value deals.",
     "Upside beyond $237K: lifting above median, plus the product-loss lever, both add to the total."])

# ==================== END ====================
out = DELIV / "CRM_Sales_Deck.pptx"
prs.save(str(out))
print("Saved:", out)